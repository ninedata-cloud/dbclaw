import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


class DocumentParser:
    """Parse various document formats to extract text content."""

    @staticmethod
    async def parse(file_path: str, file_type: str) -> Optional[str]:
        """Parse document and return text content."""
        try:
            if file_type == "txt":
                return await DocumentParser._parse_txt(file_path)
            elif file_type == "md":
                return await DocumentParser._parse_markdown(file_path)
            elif file_type == "pdf":
                return await DocumentParser._parse_pdf(file_path)
            elif file_type == "docx":
                return await DocumentParser._parse_docx(file_path)
            elif file_type == "pptx":
                return await DocumentParser._parse_pptx(file_path)
            elif file_type == "html":
                return await DocumentParser._parse_html(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
        except Exception as e:
            logger.error(f"Error parsing {file_type} file {file_path}: {e}")
            raise

    @staticmethod
    async def _parse_txt(file_path: str) -> str:
        """Parse plain text file."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    async def _parse_markdown(file_path: str) -> str:
        """Parse markdown file (preserve as-is)."""
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    async def _parse_pdf(file_path: str) -> str:
        """Parse PDF file using pypdf."""
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("pypdf is required for PDF parsing. Install with: pip install pypdf")

        reader = PdfReader(file_path)
        text_parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                text_parts.append(text)
        return "\n\n".join(text_parts)

    @staticmethod
    async def _parse_docx(file_path: str) -> str:
        """Parse DOCX file using python-docx."""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("python-docx is required for DOCX parsing. Install with: pip install python-docx")

        doc = Document(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(paragraphs)

    @staticmethod
    async def _parse_pptx(file_path: str) -> str:
        """Parse PPTX file using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX parsing. Install with: pip install python-pptx")

        prs = Presentation(file_path)
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
        return "\n\n".join(text_parts)

    @staticmethod
    async def _parse_html(file_path: str) -> str:
        """Parse HTML file using BeautifulSoup."""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 is required for HTML parsing. Install with: pip install beautifulsoup4")

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()

        soup = BeautifulSoup(html_content, "html.parser")

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Get text
        text = soup.get_text()

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = "\n".join(chunk for chunk in chunks if chunk)

        return text
